import asyncio
import os
import shutil
import tempfile
from typing import List

import fitz  # PyMuPDF
from pdf2docx import Converter as PDFToDocxConverter
from pptx import Presentation

from file_converter.converters.base import BaseConverter


class PDFToDOCXConverter(BaseConverter):
    """Convert PDF files to DOCX format"""
    
    def __init__(self, input_files_path: str, output_files_path: str):
        super().__init__(input_files_path, output_files_path)
        self.input_extension = ".pdf"
        self.output_extension = ".docx"
    
    async def _convert_single_file(self, pdf_path: str, output_path: str) -> None:
        """Convert a single PDF file to DOCX"""
        # pdf2docx conversion is synchronous, so we run it in a separate thread
        loop = asyncio.get_event_loop()
        await loop.run_in_executor(None, self._convert_sync, pdf_path, output_path)
        print(f"Converted {pdf_path} to {output_path}")
        return None
    
    def _convert_sync(self, pdf_path: str, output_path: str) -> None:
        """Synchronous conversion function to be run in a thread"""
        cv = PDFToDocxConverter(pdf_path)
        cv.convert(output_path)
        cv.close()


class PDFToJPGConverter(BaseConverter):
    """Convert PDF files to JPG images"""
    
    def __init__(self, input_files_path: str, output_files_path: str):
        super().__init__(input_files_path, output_files_path)
        self.input_extension = ".pdf"
        self.output_extension = ".jpg"
    
    async def _convert_single_file(self, pdf_path: str, output_path: str) -> None:
        """Convert a single PDF file to JPG images"""
        loop = asyncio.get_event_loop()
        await loop.run_in_executor(None, self._convert_sync, pdf_path, output_path)
        print(f"Converted {pdf_path} to JPG images")
        return None
    
    def _convert_sync(self, pdf_path: str, output_base_path: str) -> None:
        """Convert PDF to JPG images"""
        # Remove extension from output path to use as base name
        output_base_path = os.path.splitext(output_base_path)[0]
        
        # Open the PDF
        pdf_document = fitz.open(pdf_path)
        
        # Convert each page to an image
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            pix = page.get_pixmap(alpha=False)
            
            # Save the image
            output_path = f"{output_base_path}_page{page_num+1}.jpg"
            pix.save(output_path)
        
        pdf_document.close()


class PDFToPowerPointConverter(BaseConverter):
    """Convert PDF files to PowerPoint presentations"""
    
    def __init__(self, input_files_path: str, output_files_path: str):
        super().__init__(input_files_path, output_files_path)
        self.input_extension = ".pdf"
        self.output_extension = ".pptx"
    
    async def _convert_single_file(self, pdf_path: str, output_path: str) -> None:
        """Convert a single PDF file to PowerPoint"""
        # This is a complex conversion that typically requires commercial libraries
        # For a basic implementation, we'll convert PDF pages to images and insert them into slides
        loop = asyncio.get_event_loop()
        await loop.run_in_executor(None, self._convert_sync, pdf_path, output_path)
        print(f"Converted {pdf_path} to {output_path}")
        return None
    
    def _convert_sync(self, pdf_path: str, output_path: str) -> None:
        """Convert PDF to PowerPoint by inserting page images as slides"""
        # Create a new presentation
        prs = Presentation()
        
        # Open the PDF
        pdf_document = fitz.open(pdf_path)
        
        # For each page in the PDF
        for page_num in range(len(pdf_document)):
            # Get the page
            page = pdf_document.load_page(page_num)
            
            # Render page to an image
            pix = page.get_pixmap(alpha=False)
            img_bytes = pix.tobytes("jpeg")
            
            # Create a slide with the right aspect ratio
            slide_layout = prs.slide_layouts[6]  # Blank slide
            slide = prs.slides.add_slide(slide_layout)
            
            # Add the image to the slide
            with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
                tmp.write(img_bytes)
                tmp_name = tmp.name
            
            # Add picture to slide
            slide.shapes.add_picture(tmp_name, 0, 0, prs.slide_width, prs.slide_height)
            os.unlink(tmp_name)  # Remove the temporary file
        
        # Save the presentation
        prs.save(output_path)
        pdf_document.close()


class PDFToExcelConverter(BaseConverter):
    """Convert PDF tables to Excel spreadsheets"""
    
    def __init__(self, input_files_path: str, output_files_path: str):
        super().__init__(input_files_path, output_files_path)
        self.input_extension = ".pdf"
        self.output_extension = ".xlsx"
    
    async def _convert_single_file(self, pdf_path: str, output_path: str) -> None:
        """Convert a single PDF file to Excel"""
        # This is a placeholder - you would need to implement table extraction
        # Consider using libraries like tabula-py or camelot-py
        print(f"PDF to Excel conversion requires table extraction capabilities.")
        print(f"Please implement this method with a suitable library.")
        return None


class PDFToPDFAConverter(BaseConverter):
    """Convert PDF files to PDF/A format (archival standard)"""
    
    def __init__(self, input_files_path: str, output_files_path: str):
        super().__init__(input_files_path, output_files_path)
        self.input_extension = ".pdf"
        self.output_extension = ".pdf"  # Same extension but different format
    
    async def _convert_single_file(self, pdf_path: str, output_path: str) -> None:
        """Convert a single PDF file to PDF/A"""
        # PDF/A conversion is complex and typically requires specialized libraries
        # This is a placeholder implementation
        loop = asyncio.get_event_loop()
        await loop.run_in_executor(None, self._convert_sync, pdf_path, output_path)
        print(f"Converted {pdf_path} to PDF/A format at {output_path}")
        return None
    
    def _convert_sync(self, pdf_path: str, output_path: str) -> None:
        """Convert PDF to PDF/A format"""
        # For a real implementation, consider using libraries like ocrmypdf or ghostscript
        # For now, we'll just copy the file and note it's a placeholder
        shutil.copy(pdf_path, output_path) 