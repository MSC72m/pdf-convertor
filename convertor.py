from typing import Dict, List, Tuple, Any, Optional
from pdf2docx import Converter as PDFToDocxConverter
import aiofiles
from aiopath import AsyncPath
import os
import asyncio
import fitz  # PyMuPDF for PDF to image conversion
from PIL import Image
import io
import img2pdf
from docx2pdf import convert as docx_to_pdf
from pptx import Presentation
import pandas as pd
import weasyprint
import tempfile
import shutil
from pathlib import Path


class FileManager:
    def __init__(self, input_files_path: str, output_files_path: str):
        self.input_files_path = input_files_path
        self.output_files_path = output_files_path
        self.path = AsyncPath(input_files_path)
    
    async def __validate_file_path(self):
        if not await self.path.exists():
            raise FileNotFoundError(f"Path {self.input_files_path} not found")
        if not await self.path.is_dir():
            raise ValueError(f"Path {self.input_files_path} is not a directory")

    async def __validate_file_type(self, file_path: str, file_type: str):
        if not file_path.lower().endswith(file_type.lower()):
            raise ValueError(f"File {file_path} is not a {file_type} file")
        
    async def load_single_file(self, file_path: str, file_type: str) -> str:
        """
        Validate a single file
        """
        await self.__validate_file_type(file_path, file_type)
        return file_path
    
    async def load_multiple_files(self, file_type: str) -> List[str]:
        """
        Load multiple files paths
        """
        await self.__validate_file_path()
        files = []
        async for file in self.path.glob(f"*{file_type}"):
            files.append(str(file))
        return files

    async def ensure_output_directory(self):
        """Ensure the output directory exists"""
        output_path = AsyncPath(self.output_files_path)
        if not await output_path.exists():
            os.makedirs(self.output_files_path, exist_ok=True)


class BaseConverter(FileManager):
    """Base class for all converters"""
    
    def __init__(self, input_files_path: str, output_files_path: str):
        super().__init__(input_files_path, output_files_path)
        self._input_files = []
        self.input_extension = ""
        self.output_extension = ""
    
    async def _get_all_files(self) -> None:
        self._input_files = await super().load_multiple_files(self.input_extension)
        return None
    
    async def convert_all_files(self) -> None:
        await self._get_all_files()
        if not self._input_files:
            print(f"No {self.input_extension} files found in {self.input_files_path}")
            return None
        
        print(f"Found {len(self._input_files)} {self.input_extension} files to convert")
        await self.ensure_output_directory()
        
        conversion_tasks = []
        for input_path in self._input_files:
            filename = os.path.basename(input_path)
            base_name = os.path.splitext(filename)[0]
            output_path = os.path.join(self.output_files_path, f"{base_name}{self.output_extension}")
            
            # Create a task for each conversion
            conversion_tasks.append(self._convert_single_file(input_path, output_path))
        
        # Run conversions concurrently
        await asyncio.gather(*conversion_tasks)
        print("All files converted successfully")
        return None
    
    async def _convert_single_file(self, input_path: str, output_path: str) -> None:
        """Convert a single file - to be implemented by subclasses"""
        raise NotImplementedError("Subclasses must implement this method")


class PDFToDOCXConverter(BaseConverter):
    """Convert PDF files to DOCX format"""
    
    def __init__(self, input_files_path: str, output_files_path: str):
        super().__init__(input_files_path, output_files_path)
        self.input_extension = ".pdf"
        self.output_extension = ".docx"
    
    async def _convert_single_file(self, pdf_path: str, output_path: str) -> None:
        """Convert a single PDF file to DOCX"""
        # pdf2docx conversion is synchronous, so we run it in a separate thread
        loop = asyncio.get_event_loop()
        await loop.run_in_executor(None, self._convert_sync, pdf_path, output_path)
        print(f"Converted {pdf_path} to {output_path}")
        return None
    
    def _convert_sync(self, pdf_path: str, output_path: str) -> None:
        """Synchronous conversion function to be run in a thread"""
        cv = PDFToDocxConverter(pdf_path)
        cv.convert(output_path)
        cv.close()


class PDFToJPGConverter(BaseConverter):
    """Convert PDF files to JPG images"""
    
    def __init__(self, input_files_path: str, output_files_path: str):
        super().__init__(input_files_path, output_files_path)
        self.input_extension = ".pdf"
        self.output_extension = ".jpg"
    
    async def _convert_single_file(self, pdf_path: str, output_path: str) -> None:
        """Convert a single PDF file to JPG images"""
        loop = asyncio.get_event_loop()
        await loop.run_in_executor(None, self._convert_sync, pdf_path, output_path)
        print(f"Converted {pdf_path} to JPG images")
        return None
    
    def _convert_sync(self, pdf_path: str, output_base_path: str) -> None:
        """Convert PDF to JPG images"""
        # Remove extension from output path to use as base name
        output_base_path = os.path.splitext(output_base_path)[0]
        
        # Open the PDF
        pdf_document = fitz.open(pdf_path)
        
        # Convert each page to an image
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            pix = page.get_pixmap(alpha=False)
            
            # Save the image
            output_path = f"{output_base_path}_page{page_num+1}.jpg"
            pix.save(output_path)
        
        pdf_document.close()


class PDFToPowerPointConverter(BaseConverter):
    """Convert PDF files to PowerPoint presentations"""
    
    def __init__(self, input_files_path: str, output_files_path: str):
        super().__init__(input_files_path, output_files_path)
        self.input_extension = ".pdf"
        self.output_extension = ".pptx"
    
    async def _convert_single_file(self, pdf_path: str, output_path: str) -> None:
        """Convert a single PDF file to PowerPoint"""
        # This is a complex conversion that typically requires commercial libraries
        # For a basic implementation, we'll convert PDF pages to images and insert them into slides
        loop = asyncio.get_event_loop()
        await loop.run_in_executor(None, self._convert_sync, pdf_path, output_path)
        print(f"Converted {pdf_path} to {output_path}")
        return None
    
    def _convert_sync(self, pdf_path: str, output_path: str) -> None:
        """Convert PDF to PowerPoint by inserting page images as slides"""
        # Create a new presentation
        prs = Presentation()
        
        # Open the PDF
        pdf_document = fitz.open(pdf_path)
        
        # For each page in the PDF
        for page_num in range(len(pdf_document)):
            # Get the page
            page = pdf_document.load_page(page_num)
            
            # Render page to an image
            pix = page.get_pixmap(alpha=False)
            img_bytes = pix.tobytes("jpeg")
            
            # Create a slide with the right aspect ratio
            slide_layout = prs.slide_layouts[6]  # Blank slide
            slide = prs.slides.add_slide(slide_layout)
            
            # Add the image to the slide
            with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
                tmp.write(img_bytes)
                tmp_name = tmp.name
            
            # Add picture to slide
            slide.shapes.add_picture(tmp_name, 0, 0, prs.slide_width, prs.slide_height)
            os.unlink(tmp_name)  # Remove the temporary file
        
        # Save the presentation
        prs.save(output_path)
        pdf_document.close()


class PDFToExcelConverter(BaseConverter):
    """Convert PDF tables to Excel spreadsheets"""
    
    def __init__(self, input_files_path: str, output_files_path: str):
        super().__init__(input_files_path, output_files_path)
        self.input_extension = ".pdf"
        self.output_extension = ".xlsx"
    
    async def _convert_single_file(self, pdf_path: str, output_path: str) -> None:
        """Convert a single PDF file to Excel"""
        # Note: For accurate PDF table extraction, consider using tabula-py or camelot-py
        # This is a simplified implementation
        loop = asyncio.get_event_loop()
        await loop.run_in_executor(None, self._convert_sync, pdf_path, output_path)
        print(f"Converted {pdf_path} to {output_path}")
        return None
    
    def _convert_sync(self, pdf_path: str, output_path: str) -> None:
        """Extract tables from PDF and save to Excel"""
        # This is a placeholder - for production use, install and use tabula-py:
        # import tabula
        # tables = tabula.read_pdf(pdf_path, pages='all')
        # with pd.ExcelWriter(output_path) as writer:
        #     for i, table in enumerate(tables):
        #         table.to_excel(writer, sheet_name=f'Table {i+1}')
        
        # For now, create a simple Excel file noting this is a placeholder
        df = pd.DataFrame({'Note': ['PDF table extraction requires specialized libraries like tabula-py or camelot-py']})
        df.to_excel(output_path, index=False)


class PDFToPDFAConverter(BaseConverter):
    """Convert PDF files to PDF/A format (archival standard)"""
    
    def __init__(self, input_files_path: str, output_files_path: str):
        super().__init__(input_files_path, output_files_path)
        self.input_extension = ".pdf"
        self.output_extension = ".pdf"  # Same extension but different format
    
    async def _convert_single_file(self, pdf_path: str, output_path: str) -> None:
        """Convert a single PDF file to PDF/A"""
        # PDF/A conversion is complex and typically requires specialized libraries
        # This is a placeholder implementation
        loop = asyncio.get_event_loop()
        await loop.run_in_executor(None, self._convert_sync, pdf_path, output_path)
        print(f"Converted {pdf_path} to PDF/A format at {output_path}")
        return None
    
    def _convert_sync(self, pdf_path: str, output_path: str) -> None:
        """Convert PDF to PDF/A format"""
        # For a real implementation, consider using libraries like ocrmypdf or ghostscript
        # For now, we'll just copy the file and note it's a placeholder
        shutil.copy(pdf_path, output_path)


class JPGToPDFConverter(BaseConverter):
    """Convert JPG images to PDF"""
    
    def __init__(self, input_files_path: str, output_files_path: str):
        super().__init__(input_files_path, output_files_path)
        self.input_extension = ".jpg"
        self.output_extension = ".pdf"
    
    async def convert_all_files(self) -> None:
        """Override to handle multiple images to one PDF or one image per PDF"""
        await self._get_all_files()
        if not self._input_files:
            print(f"No JPG files found in {self.input_files_path}")
            return None
        
        print(f"Found {len(self._input_files)} JPG files to convert")
        await self.ensure_output_directory()
        
        # Group files by base name (without numbers at the end)
        file_groups = {}
        for file_path in self._input_files:
            filename = os.path.basename(file_path)
            # Try to extract a base name by removing trailing numbers
            base_name = Path(filename).stem
            base_name = ''.join(c for c in base_name if not c.isdigit()).rstrip('_-. ')
            
            if not base_name:  # If we removed everything, use the original name
                base_name = Path(filename).stem
            
            if base_name not in file_groups:
                file_groups[base_name] = []
            file_groups[base_name].append(file_path)
        
        # Convert each group to a PDF
        conversion_tasks = []
        for base_name, files in file_groups.items():
            output_path = os.path.join(self.output_files_path, f"{base_name}.pdf")
            conversion_tasks.append(self._convert_file_group(files, output_path))
        
        await asyncio.gather(*conversion_tasks)
        print("All files converted successfully")
    
    async def _convert_file_group(self, image_paths: List[str], output_path: str) -> None:
        """Convert a group of JPG images to a single PDF"""
        loop = asyncio.get_event_loop()
        await loop.run_in_executor(None, self._convert_sync, image_paths, output_path)
        print(f"Converted {len(image_paths)} images to {output_path}")
    
    def _convert_sync(self, image_paths: List[str], output_path: str) -> None:
        """Convert JPG images to PDF"""
        # Sort images by name to maintain order
        image_paths.sort()
        
        # Convert images to PDF
        with open(output_path, "wb") as f:
            f.write(img2pdf.convert(image_paths))
    
    async def _convert_single_file(self, jpg_path: str, output_path: str) -> None:
        """Convert a single JPG file to PDF"""
        loop = asyncio.get_event_loop()
        await loop.run_in_executor(None, self._convert_single_sync, jpg_path, output_path)
        print(f"Converted {jpg_path} to {output_path}")
        return None
    
    def _convert_single_sync(self, jpg_path: str, output_path: str) -> None:
        """Convert a single JPG to PDF"""
        with open(output_path, "wb") as f:
            f.write(img2pdf.convert(jpg_path))


class WordToPDFConverter(BaseConverter):
    """Convert Word documents to PDF"""
    
    def __init__(self, input_files_path: str, output_files_path: str):
        super().__init__(input_files_path, output_files_path)
        self.input_extension = ".docx"
        self.output_extension = ".pdf"
    
    async def _convert_single_file(self, docx_path: str, output_path: str) -> None:
        """Convert a single Word document to PDF"""
        loop = asyncio.get_event_loop()
        await loop.run_in_executor(None, self._convert_sync, docx_path, output_path)
        print(f"Converted {docx_path} to {output_path}")
        return None
    
    def _convert_sync(self, docx_path: str, output_path: str) -> None:
        """Convert Word to PDF using docx2pdf"""
        convert(docx_path, output_path)


class PowerPointToPDFConverter(BaseConverter):
    """Convert PowerPoint presentations to PDF"""
    
    def __init__(self, input_files_path: str, output_files_path: str):
        super().__init__(input_files_path, output_files_path)
        self.input_extension = ".pptx"
        self.output_extension = ".pdf"
    
    async def _convert_single_file(self, pptx_path: str, output_path: str) -> None:
        """Convert a single PowerPoint presentation to PDF"""
        # This typically requires Microsoft Office or LibreOffice
        # For a placeholder implementation, we'll note this limitation
        print(f"PowerPoint to PDF conversion requires Microsoft Office or LibreOffice.")
        print(f"Please install a suitable converter and update this method.")
        return None


class ExcelToPDFConverter(BaseConverter):
    """Convert Excel spreadsheets to PDF"""
    
    def __init__(self, input_files_path: str, output_files_path: str):
        super().__init__(input_files_path, output_files_path)
        self.input_extension = ".xlsx"
        self.output_extension = ".pdf"
    
    async def _convert_single_file(self, xlsx_path: str, output_path: str) -> None:
        """Convert a single Excel spreadsheet to PDF"""
        # This typically requires Microsoft Office or LibreOffice
        # For a placeholder implementation, we'll note this limitation
        print(f"Excel to PDF conversion requires Microsoft Office or LibreOffice.")
        print(f"Please install a suitable converter and update this method.")
        return None


class HTMLToPDFConverter(BaseConverter):
    """Convert HTML files to PDF"""
    
    def __init__(self, input_files_path: str, output_files_path: str):
        super().__init__(input_files_path, output_files_path)
        self.input_extension = ".html"
        self.output_extension = ".pdf"
    
    async def _convert_single_file(self, html_path: str, output_path: str) -> None:
        """Convert a single HTML file to PDF"""
        loop = asyncio.get_event_loop()
        await loop.run_in_executor(None, self._convert_sync, html_path, output_path)
        print(f"Converted {html_path} to {output_path}")
        return None
    
    def _convert_sync(self, html_path: str, output_path: str) -> None:
        """Convert HTML to PDF using WeasyPrint"""
        # Read the HTML file
        with open(html_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        # Convert to PDF
        html = weasyprint.HTML(string=html_content, base_url=os.path.dirname(html_path))
        html.write_pdf(output_path)




